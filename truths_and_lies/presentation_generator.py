from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from .models import ListOfParticipants
from random import shuffle


PPTX_TEMPLATE = "ppt-templates/funny-presentation.pptx"
FONT_SIZE = Pt(25)


def create_presentation(
    participants: ListOfParticipants,
    presentation_path: Path,
) -> None:
    if len(participants.participants) == 0:
        raise ValueError("No participants found.")

    prs = Presentation(PPTX_TEMPLATE)

    # Change the title of the first slide

    slide = prs.slides[0]
    title = slide.shapes.title
    title.text = f"Truths and Lies"
    subtitle = slide.placeholders[1]
    subtitle.text = "Fiesta Edy 2023"

    for participant in participants.participants:
        # Add a big slide with the participant name
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = participant.name.capitalize()
        subtitle = slide.placeholders[1]
        subtitle.text = "Te toca a ti!!"

        # Shuffle the statements
        statements = participant.statements.copy()
        shuffle(statements)

        # Add incremental slides for each statement
        for i in range(1, len(statements) + 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"Verdades de: {participant.name.capitalize()}"
            content = slide.placeholders[1]
            for statement in statements[:i]:
                p = content.text_frame.add_paragraph()
                p.text = statement.statement
                p.font.size = FONT_SIZE

        # Add a final slide with all statements colored
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"Momento de la verdad: {participant.name.capitalize()}"
        content = slide.placeholders[1]
        content.text = ""
        for statement in statements:
            color = RGBColor(0, 128, 0) if not statement.isLie else RGBColor(255, 0, 0)
            p = content.text_frame.add_paragraph()
            p.text = statement.statement
            p.font.color.rgb = color
            p.font.size = FONT_SIZE

    prs.save(presentation_path)
